from pptx import Presentation

def extract_slides_with_asterisk(pptx_path, output_text_file):
    # loads the presentation
    presentation = Presentation(pptx_path)

    # opens the output text file
    with open(output_text_file, 'w', encoding='utf-8') as file:
        # iterates through the slides
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = ''
            # extracts text from all shapes on the slide
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + '\n'
            
            # checks if the slide contains an asterisk
            if '*' in slide_text:
                # writes the slide number and its text to the file
                file.write(f"Slide {slide_num}:\n")
                file.write(slide_text)
                file.write("\n" + "="*40 + "\n")

if __name__ == "__main__":
    pptx_path = "GEP1018.pptx"
    output_text_file = "GEPslides.txt"
    extract_slides_with_asterisk(pptx_path, output_text_file)
